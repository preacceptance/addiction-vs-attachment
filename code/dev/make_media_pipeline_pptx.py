#!/usr/bin/env python3
"""
Figure 2. Media text analysis workflow — black-and-white pptx.

Parallel to the legal `make_pipeline_pptx.py` / figure1_pipeline.pptx, and
matching the edited final style of that diagram:
  - stage tabs "Preparation / Surface Analysis / Deeper Analysis / Output"
  - "mentions of" (surface) vs "meanings of" (deeper) parallel wording
  - enumerated (i)/(ii)/(iii)/(iv) category list
  - "Surface analysis result" side callout with live counts
  - genericized "Validation" callout (no hardcoded kappa)

Media-specific: the unit of analysis is the PARAGRAPH, extracted from articles
that are themselves split out of Factiva PDF search exports (blocks mode).
Counts are unique-usable (full articles, cross-search duplicates removed).

Numbers pulled live from media_articles.xlsx + media_paragraphs.xlsx.

Output: figure2_media_pipeline.pptx
"""
from __future__ import annotations
from pathlib import Path
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ── Live numbers ─────────────────────────────────────────────────────────
arts = pd.read_excel("media_articles.xlsx")
pars = pd.read_excel("media_paragraphs.xlsx")

N_PDFS        = arts["pdf"].nunique()
N_ART_FULL    = int(arts["usable"].sum())
N_ART_UNIQUE  = int((arts["usable"] & ~arts["is_duplicate"]).sum())

par_u = pars[pars["article_usable"] & ~pars["is_duplicate"]]
N_PARAS   = len(par_u)
N_ADDICT  = int((par_u["surface_meaning"] == "addiction").sum())
N_ATTACH  = int((par_u["surface_meaning"] == "attachment").sum())
N_BOTH    = int((par_u["surface_meaning"] == "both").sum())
N_NEITHER = int((par_u["surface_meaning"] == "neither").sum())

# ── Slide setup ──────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank


# ── Helpers ──────────────────────────────────────────────────────────────
def add_text_shape(left, top, width, height, lines, *,
                   size_first=12, size_body=10, bold_first=True,
                   align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    shp.fill.solid(); shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = BLACK; shp.line.width = Pt(1.0)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.10)
    tf.margin_top = tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        for run in p.runs:
            run.font.name = "Calibri"
            run.font.color.rgb = BLACK
            run.font.size = Pt(size_first if (i == 0 and bold_first) else size_body)
            run.font.bold = (i == 0 and bold_first)
    return shp


def add_tab(left, top, width, height, label):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    shp.fill.solid(); shp.fill.fore_color.rgb = BLACK
    shp.line.color.rgb = BLACK; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = label; p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.name = "Calibri"; run.font.size = Pt(13)
        run.font.bold = True; run.font.color.rgb = WHITE


def add_arrow(start, end):
    sx, sy = start; ex, ey = end
    conn = slide.shapes.add_connector(2, Inches(sx), Inches(sy),
                                      Inches(ex), Inches(ey))
    conn.line.color.rgb = BLACK; conn.line.width = Pt(1.5)
    from pptx.oxml.ns import qn
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"),
                             {"type": "triangle", "w": "med", "len": "med"}))


# ── Title ────────────────────────────────────────────────────────────────
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.16), Inches(9.0), Inches(0.45))
tf = title.text_frame; tf.margin_top = tf.margin_bottom = Inches(0)
p = tf.paragraphs[0]; p.text = "Figure 2. Media text analysis workflow"; p.alignment = PP_ALIGN.LEFT
for run in p.runs:
    run.font.name = "Calibri"; run.font.size = Pt(18); run.font.bold = True
    run.font.color.rgb = BLACK


# ── Layout (matches the edited legal figure) ─────────────────────────────
TAB_X, TAB_W = 0.22, 1.23
MAIN_X, MAIN_W = 1.64, 4.85
SIDE_X, SIDE_W = 6.82, 3.10

ROW_PREP_Y, ROW_PREP_H = 0.95, 1.30
ROW_SURF_Y, ROW_SURF_H = 2.55, 1.30
ROW_DEEP_Y, ROW_DEEP_H = 4.15, 1.40
ROW_OUT_Y,  ROW_OUT_H  = 5.85, 1.00

def tab_y(row_y, row_h, tab_h=0.60):
    return row_y + row_h / 2 - tab_h / 2

add_tab(TAB_X, tab_y(ROW_PREP_Y, ROW_PREP_H), TAB_W, 0.60, "Preparation")
add_tab(TAB_X, tab_y(ROW_SURF_Y, ROW_SURF_H), TAB_W, 0.60, "Surface Analysis")
add_tab(TAB_X, tab_y(ROW_DEEP_Y, ROW_DEEP_H), TAB_W, 0.60, "Deeper Analysis")
add_tab(TAB_X, tab_y(ROW_OUT_Y,  ROW_OUT_H),  TAB_W, 0.58, "Output")


# ── Main content boxes ───────────────────────────────────────────────────
add_text_shape(MAIN_X, ROW_PREP_Y, MAIN_W, ROW_PREP_H, [
    "Data preparation",
    f"{N_PDFS} Factiva search exports split into articles, then numbered paragraphs.",
    "Preview-only stubs and cross-search duplicate articles flagged out.",
    f"Corpus: n = {N_ART_UNIQUE:,} unique articles → {N_PARAS:,} paragraphs.",
])

add_text_shape(MAIN_X, ROW_SURF_Y, MAIN_W, ROW_SURF_H, [
    "Surface analysis",
    "Dictionary-based tagging on word stems of 'attachment' and 'addiction'.",
    "Each paragraph classified as having mentions of: ",
    "(i) addiction, (ii) attachment, (iii) both, or (iv) neither.",
])

add_text_shape(MAIN_X, ROW_DEEP_Y, MAIN_W, ROW_DEEP_H, [
    "Deeper analysis",
    "GPT-5 four-class classifier with 30 few-shot examples.",
    "Each paragraph classified as having meanings of: ",
    "(i) addiction, (ii) attachment, (iii) both, or (iv) neither.",
])

add_text_shape(MAIN_X, ROW_OUT_Y, MAIN_W, ROW_OUT_H, [
    "Final labeled corpus",
    f"n = {N_PARAS:,} paragraphs across four categories.",
])


# ── Side callouts ────────────────────────────────────────────────────────
side_y = ROW_SURF_Y + (ROW_SURF_H - 1.16) / 2
add_text_shape(SIDE_X, side_y, SIDE_W, 1.16, [
    "Surface analysis result",
    f"Addiction   n = {N_ADDICT}",
    f"Attachment  n = {N_ATTACH}",
    f"Both        n = {N_BOTH}",
    f"Neither     n = {N_NEITHER:,}",
], size_first=11, size_body=10)
add_arrow((MAIN_X + MAIN_W, side_y + 0.58), (SIDE_X, side_y + 0.58))

side_y2 = ROW_DEEP_Y + (ROW_DEEP_H - 1.21) / 2
add_text_shape(SIDE_X, side_y2, SIDE_W, 1.21, [
    "Validation",
    "Compared against manual coding of subset (n = 30 paragraphs).",
    "Inter-rater reliability: Cohen's and Fleiss' kappa.",
], size_first=11, size_body=10)
add_arrow((MAIN_X + MAIN_W, side_y2 + 0.60), (SIDE_X, side_y2 + 0.60))


# ── Vertical arrows between main boxes ────────────────────────────────────
ax = MAIN_X + MAIN_W / 2
add_arrow((ax, ROW_PREP_Y + ROW_PREP_H), (ax, ROW_SURF_Y))
add_arrow((ax, ROW_SURF_Y + ROW_SURF_H), (ax, ROW_DEEP_Y))
add_arrow((ax, ROW_DEEP_Y + ROW_DEEP_H), (ax, ROW_OUT_Y))


out = Path("figure2_media_pipeline.pptx")
prs.save(out)
print(f"wrote {out}")
print(f"  {N_PDFS} PDFs | {N_ART_UNIQUE} unique articles | {N_PARAS:,} paragraphs")
print(f"  surface: addiction={N_ADDICT} attachment={N_ATTACH} both={N_BOTH} neither={N_NEITHER:,}")
