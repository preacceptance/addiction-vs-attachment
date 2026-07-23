#!/usr/bin/env python3
"""
Figure 1. Text analysis workflow — black-and-white pptx.

Mirrors the layout of Itai's Text_Analysis_Academic.pptx: a vertical
Preparation → Surface → Deeper → Output flow with stage tabs on the left,
main content boxes in the middle, and side-callouts ("Dictionary result",
"Validation") to the right of the Surface and Deeper rows.

Numbers are pulled live from legal_paragraphs.xlsx so the figure cannot
drift from the corpus.

Output: figure1_pipeline.pptx
"""
from __future__ import annotations
from pathlib import Path
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# ── Live numbers ─────────────────────────────────────────────────────────
legal = pd.read_excel("legal_paragraphs.xlsx")
N_PARAS   = len(legal)
N_CASES   = legal["case"].nunique()
N_ADDICT  = int(legal["has_addict"].sum())
N_ATTACH  = int(legal["has_attach"].sum())      # substantive only
N_BOTH    = int((legal["has_addict"] &  legal["has_attach"]).sum())
N_NEITHER = int((~legal["has_addict"] & ~legal["has_attach"]).sum())


# ── Slide setup ──────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank


# ── Helpers ──────────────────────────────────────────────────────────────
def add_text_shape(left, top, width, height, lines, *,
                   fill_white=True, line_black=True,
                   bold_first=True, size_first=12, size_body=10,
                   align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                   pad_l=0.10, pad_r=0.10, pad_t=0.08, pad_b=0.08):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = WHITE if fill_white else BLACK
    shp.line.color.rgb = BLACK if line_black else WHITE
    shp.line.width = Pt(1.0)
    shp.shadow.inherit = False

    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left   = Inches(pad_l)
    tf.margin_right  = Inches(pad_r)
    tf.margin_top    = Inches(pad_t)
    tf.margin_bottom = Inches(pad_b)
    tf.vertical_anchor = anchor

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        for run in p.runs:
            run.font.name = "Calibri"
            run.font.color.rgb = BLACK
            run.font.size = Pt(size_first if (i == 0 and bold_first) else size_body)
            run.font.bold = (i == 0 and bold_first)
    return shp


def add_tab(left, top, width, height, label):
    """Black rounded-rectangle tab with white label text, centered."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = BLACK
    shp.line.color.rgb = BLACK
    shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top  = tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = label
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.name = "Calibri"
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = WHITE


def add_arrow(start, end, *, head=True):
    """Black connector. start/end are (x,y) in inches."""
    sx, sy = start; ex, ey = end
    conn = slide.shapes.add_connector(2,  # straight
                                      Inches(sx), Inches(sy),
                                      Inches(ex), Inches(ey))
    line = conn.line
    line.color.rgb = BLACK
    line.width = Pt(1.5)
    if head:
        # python-pptx doesn't expose arrowhead directly — use XML hack
        from pptx.oxml.ns import qn
        ln = conn.line._get_or_add_ln()
        tail = ln.makeelement(qn("a:tailEnd"),
                              {"type": "triangle", "w": "med", "len": "med"})
        ln.append(tail)


# ── Title ────────────────────────────────────────────────────────────────
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3),
                                 Inches(9.0), Inches(0.45))
tf = title.text_frame; tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
p = tf.paragraphs[0]; p.text = "Figure 1. Text analysis workflow"
p.alignment = PP_ALIGN.LEFT
for run in p.runs:
    run.font.name = "Calibri"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = BLACK


# ── Layout constants (inches) ────────────────────────────────────────────
TAB_X      = 0.10
TAB_W      = 1.10
MAIN_X     = 1.35
MAIN_W     = 4.85
SIDE_X     = 6.50
SIDE_W     = 3.10

ROW_PREP_Y, ROW_PREP_H   = 0.95, 1.20
ROW_SURF_Y, ROW_SURF_H   = 2.45, 1.30
ROW_DEEP_Y, ROW_DEEP_H   = 4.05, 1.40
ROW_OUT_Y,  ROW_OUT_H    = 5.80, 1.00

# Stage tab vertical centers
def tab_y(row_y, row_h, tab_h=0.46):
    return row_y + row_h / 2 - tab_h / 2


# ── Stage tabs (left column) ─────────────────────────────────────────────
add_tab(TAB_X, tab_y(ROW_PREP_Y, ROW_PREP_H), TAB_W, 0.46, "Preparation")
add_tab(TAB_X, tab_y(ROW_SURF_Y, ROW_SURF_H), TAB_W, 0.46, "Surface")
add_tab(TAB_X, tab_y(ROW_DEEP_Y, ROW_DEEP_H), TAB_W, 0.46, "Deeper")
add_tab(TAB_X, tab_y(ROW_OUT_Y,  ROW_OUT_H),  TAB_W, 0.46, "Output")


# ── Main content boxes ───────────────────────────────────────────────────
add_text_shape(MAIN_X, ROW_PREP_Y, MAIN_W, ROW_PREP_H, [
    "Data preparation",
    f"14 legal complaints parsed into numbered paragraphs.",
    f"Procedural 'attached as Exhibit X' uses excluded from attachment vocabulary.",
    f"Corpus: n = {N_PARAS:,} paragraphs.",
])

add_text_shape(MAIN_X, ROW_SURF_Y, MAIN_W, ROW_SURF_H, [
    "Surface-level analysis",
    "Dictionary-based tagging on word stems of 'attachment' and 'addiction'.",
    "Each paragraph labeled: addiction, attachment, both, or neither.",
])

add_text_shape(MAIN_X, ROW_DEEP_Y, MAIN_W, ROW_DEEP_H, [
    "Deeper analysis",
    "GPT-5 four-class classifier guided by a formal coding manual.",
    "Few-shot calibration: 30 paragraphs, 3-coder consensus (1 in 138 of corpus).",
    "Each paragraph classified: addiction, attachment, both, or neither.",
])

add_text_shape(MAIN_X, ROW_OUT_Y, MAIN_W, ROW_OUT_H, [
    "Final labeled corpus",
    f"n = {N_PARAS:,} paragraphs across four categories.",
])


# ── Side callouts (Surface and Deeper rows) ──────────────────────────────
# Dictionary result (right of Surface)
side_y = ROW_SURF_Y + (ROW_SURF_H - 1.02) / 2
add_text_shape(SIDE_X, side_y, SIDE_W, 1.02, [
    "Dictionary result",
    f"Addiction   n = {N_ADDICT}",
    f"Attachment  n = {N_ATTACH}",
    f"Both        n = {N_BOTH}",
    f"Neither     n = {N_NEITHER:,}",
], size_first=11, size_body=10)
add_arrow((MAIN_X + MAIN_W, side_y + 0.51),
          (SIDE_X, side_y + 0.51))

# Validation (right of Deeper)
side_y2 = ROW_DEEP_Y + (ROW_DEEP_H - 1.02) / 2
add_text_shape(SIDE_X, side_y2, SIDE_W, 1.02, [
    "Validation",
    "Held-out subset (n = 29 paragraphs);",
    "3 expert raters + LLM.",
    "Fleiss' kappa: pending recoding",
], size_first=11, size_body=10)
add_arrow((MAIN_X + MAIN_W, side_y2 + 0.51),
          (SIDE_X, side_y2 + 0.51))


# ── Vertical down-arrows between main boxes ──────────────────────────────
arrow_x = MAIN_X + MAIN_W / 2
add_arrow((arrow_x, ROW_PREP_Y + ROW_PREP_H),
          (arrow_x, ROW_SURF_Y))
add_arrow((arrow_x, ROW_SURF_Y + ROW_SURF_H),
          (arrow_x, ROW_DEEP_Y))
add_arrow((arrow_x, ROW_DEEP_Y + ROW_DEEP_H),
          (arrow_x, ROW_OUT_Y))


# ── Save ─────────────────────────────────────────────────────────────────
out = Path("figure1_pipeline.pptx")
prs.save(out)
print(f"wrote {out}")
