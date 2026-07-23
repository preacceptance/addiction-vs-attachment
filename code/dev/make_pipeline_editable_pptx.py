#!/usr/bin/env python3
"""
Build figures/figure1_pipeline_editable.pptx — the text-analysis workflow as an
editable PowerPoint: a left rail of stage labels, four stacked content boxes, and
down-arrows between them. NO side offshoots. Every box is a real text box the user
can edit directly in PowerPoint. Numbers/wording match the manuscript's Figure 1.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent.parent
OUT  = ROOT / "figures" / "figure1_pipeline_editable.pptx"

NAVY  = RGBColor(0x3A, 0x4A, 0x6B)
LIGHT = RGBColor(0xF4, 0xF6, 0xFB)
DARK  = RGBColor(0x1F, 0x2A, 0x3D)
BODY  = RGBColor(0x22, 0x22, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank


def _no_border(shape):
    shape.line.fill.background()


def title_box(x, y, w, h, text, size, color, bold=True, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    r.font.name = "Arial"
    return tb


def rail_chip(y_center, label):
    h, w, x = 0.72, 1.95, 0.5
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y_center - h/2), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = NAVY; _no_border(sp)
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Arial"


def content_box(top, h, header, lines):
    x, w = 2.75, 9.6
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(top), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = LIGHT
    sp.line.color.rgb = NAVY; sp.line.width = Pt(1.4)
    tf = sp.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.22); tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = header
    r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = DARK; r.font.name = "Arial"
    for ln in lines:
        bp = tf.add_paragraph()
        br = bp.add_run(); br.text = ln
        br.font.size = Pt(12); br.font.color.rgb = BODY; br.font.name = "Arial"
    return x + w/2


def down_arrow(x_center, y_top, y_bottom):
    w = 0.34
    sp = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                Inches(x_center - w/2), Inches(y_top),
                                Inches(w), Inches(y_bottom - y_top))
    sp.fill.solid(); sp.fill.fore_color.rgb = NAVY; _no_border(sp)


title_box(0.5, 0.15, 12.3, 0.6, "Figure 1. Text analysis workflow", 24, DARK)

stages = [
    ("Preparation", "Data preparation", [
        "Legal: 14 complaints → 4,146 numbered paragraphs.",
        "Media: 828 articles → 578 unique → 21,314 paragraphs.",
        "Procedural “attached as Exhibit” uses excluded.",
    ]),
    ("Surface\nanalysis", "Surface-level analysis", [
        "Dictionary tagging on the word stems of “addiction” and “attachment.”",
        "Each paragraph labeled addiction, attachment, both, or neither.",
    ]),
    ("Deeper\nanalysis", "Deeper analysis", [
        "GPT-5.4 four-class classifier (high reasoning), guided by a formal",
        "coding manual, one per corpus.",
        "Calibrated with 150 example paragraphs per corpus.",
    ]),
    ("Output", "Final labeled corpus", [
        "Legal: 4,146 paragraphs.   Media: 21,314 paragraphs.",
        "Each labeled addiction, attachment, both, or neither.",
    ]),
]

tops    = [0.95, 2.55, 4.15, 5.75]
heights = [1.35, 1.15, 1.35, 1.05]
for i, (rail, head, body) in enumerate(stages):
    top, h = tops[i], heights[i]
    rail_chip(top + h/2, rail)
    xc = content_box(top, h, head, body)
    if i > 0:
        prev_bottom = tops[i-1] + heights[i-1]
        down_arrow(xc, prev_bottom + 0.02, top - 0.02)

prs.save(OUT)
print(f"wrote {OUT}")
